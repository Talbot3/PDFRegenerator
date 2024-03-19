import os

import  fitz  # PyMuPDF
from reportlab.pdfgen import canvas
from PyPDF2 import PdfReader, PdfWriter
import io
import click

from pptx import Presentation
def ppt_to_pdf(input_ppt_file, output_pdf_file):
    prs = Presentation(input_ppt_file)
    pdf_canvas = canvas.Canvas(output_pdf_file)

    for slide in prs.slides:
        pdf_canvas.drawString(100, 800, slide.title.text)  # 示例：将PPT每页的标题绘制到PDF页面上

        # 如果要将PPT页面内容精确转换为PDF页面内容，需要对slide.shapes进行遍历，根据需要重新排版

        pdf_canvas.showPage()

    pdf_canvas.save()

def ppt_dir_to_pdf(input_ppt_dir, output_pdf_dir):
    import os

    if not os.path.exists(output_pdf_dir):
        os.makedirs(output_pdf_dir)

    for filename in os.listdir(input_ppt_dir):
        if filename.endswith(".pptx"):  # 确保只处理PPTX文件
            ppt_file = os.path.join(input_ppt_dir, filename)
            pdf_file = os.path.join(output_pdf_dir, filename.replace(".pptx", ".pdf"))
            ppt_to_pdf(ppt_file, pdf_file)


@click.command()
@click.argument('pdf_path', type=click.Path())
@click.option('output_path', '-o', help='Output file path', type=click.Path(), default=lambda: os.path.join(os.getcwd()))
def create_pdf_with_contents(pdf_path, output_path):
    doc = fitz.open(os.path.abspath(pdf_path))
    toc = []
    titles = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text").split('\n')[0].strip()
        if text:
            toc.append((page_num, text))
            titles.append(text)

    packet = io.BytesIO()
    c = canvas.Canvas(packet)
    for i, title in enumerate(titles):
        print(title)
        c.drawString(72, 800 - i * 10, f"{title} ...... Page {toc[i][0] + 1}")
    c.save()

    packet.seek(0)
    new_pdf = PdfReader(packet) # 更新类名
    output = PdfWriter() # 更新类名
    output.add_page(new_pdf.pages[0]) # 更新调用方式

    existing_pdf = PdfReader(open(pdf_path, "rb")) # 更新类名
    for i in range(len(existing_pdf.pages)):
        output.add_page(existing_pdf.pages[i]) # 更新调用方式

    file_name, file_extension = os.path.splitext(os.path.basename(pdf_path))
    # 添加标记到文件名
    if output_path:
        mark = 'mark'
        output_file_path = os.path.join(os.path.dirname(os.path.abspath(output_path)), f'{file_name}_{mark}{file_extension}')

    with open(output_file_path, "wb") as outputStream:
        print(f"regenerator path -> {output_file_path}")
        output.write(outputStream)
if __name__ == "__main__":
    create_pdf_with_contents()
